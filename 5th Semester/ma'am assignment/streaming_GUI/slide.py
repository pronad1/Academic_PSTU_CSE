from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Helper to create a standard slide
    def add_slide(title_text, content_points):
        slide_layout = prs.slide_layouts[1]  # Bullet point layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        
        # Content
        tf = slide.placeholders[1].text_frame
        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(20)

    # --- SLIDE 1: Title ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "UDP-Based Media Streaming Application"
    subtitle.text = "Final Lab Project Implementation\nConnectionless Sockets & Random Packet Distribution\n\nPresented by: [Your Name]"

    # --- SLIDE 2: Objective ---
    add_slide("Project Objective", [
        "Develop a client-server multimedia streaming application.",
        "Utilize Connectionless Sockets (UDP) instead of TCP.",
        "Implement random data packet sizes (1000-2000 bytes).",
        "Achieve Progressive Playback (Watch while downloading).",
        "Create a user-friendly GUI for easy operation."
    ])

    # --- SLIDE 3: Requirements ---
    add_slide("Lab Requirements (The Problem)", [
        "1. Protocol: Application must use UDP datagrams.",
        "2. Server: Read file chunks randomly between 1000-2000 bytes.",
        "3. Variable Packets: Handle last packet size < 1000 bytes.",
        "4. Client: Receive bytes and reconstruct the media file.",
        "5. Playback: Launch media player automatically once buffer is filled."
    ])

    # --- SLIDE 4: Theory ---
    add_slide("Why UDP for Streaming?", [
        "TCP (Transmission Control Protocol):",
        " - Reliable but slow (Handshakes, ACKs, Retransmission).",
        " - High latency is bad for real-time video.",
        "UDP (User Datagram Protocol):",
        " - Connectionless (No handshake).",
        " - 'Best Effort' delivery (Fast).",
        " - Ideal for streaming where speed > 100% accuracy."
    ])

    # --- SLIDE 5: Architecture ---
    add_slide("System Architecture", [
        "1. Request: Client sends 'START' command to Server.",
        "2. Processing: Server reads file in random chunks.",
        "3. Transport: Chunks sent as UDP Datagrams over network.",
        "4. Assembly: Client writes received bytes to temp file.",
        "5. Trigger: Media player launches when buffer hits 100KB.",
        "(Visual diagram to be explained during demo)"
    ])

    # --- SLIDE 6: Code Logic (Server) ---
    add_slide("Implementation: Server Side", [
        "Meeting the Random Size Requirement:",
        "",
        "chunk_size = random.randint(1000, 2000)",
        "chunk = file.read(chunk_size)",
        "server_socket.sendto(chunk, client_address)",
        "",
        "- Uses Python 'random' library.",
        "- Dynamically changes size for every packet.",
        "- Handles End-of-File naturally."
    ])

    # --- SLIDE 7: Code Logic (Client) ---
    add_slide("Implementation: Client Side", [
        "Meeting the Playback Requirement:",
        "- Multi-threaded approach (GUI + Receiver Thread).",
        "",
        "if bytes_received >= BUFFER_THRESHOLD:",
        "    if not playback_started:",
        "        launch_media_player()",
        "",
        "- Ensures video plays before download completes.",
        "- Keeps UI responsive during transfer."
    ])

    # --- SLIDE 8: GUI Features ---
    add_slide("Application Features", [
        "Built with Python Tkinter.",
        "Dual Functionality: Acts as both Client and Server.",
        "Real-time Statistics:",
        " - Download Speed (KB/s)",
        " - Progress Bar & Percentage",
        " - ETA Calculator",
        "Robust Error Handling & Logging."
    ])

    # --- SLIDE 9: DEMO ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Live Demonstration"
    tf = slide.placeholders[1].text_frame
    tf.text = "Switching to Application View..."
    p = tf.add_paragraph()
    p.text = "1. Server Setup & File Selection"
    p = tf.add_paragraph()
    p.text = "2. Client Connection"
    p = tf.add_paragraph()
    p.text = "3. Streaming & Playback Verification"

    # --- SLIDE 10: Conclusion ---
    add_slide("Conclusion", [
        "Summary:",
        "- Successfully implemented connectionless streaming.",
        "- Met all strict packet size requirements.",
        "",
        "Future Improvements:",
        "- Add Sequence Numbers for packet ordering.",
        "- Implement reliability layer (ACKs) if needed.",
        "",
        "Thank You. Questions?"
    ])

    prs.save('Final_Project_Presentation.pptx')
    print("Presentation saved as 'Final_Project_Presentation.pptx'")

if __name__ == "__main__":
    create_presentation()